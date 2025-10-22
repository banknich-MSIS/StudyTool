"""
File processing utilities for extracting text from various file formats.
Supports PDF, DOCX, PPTX, and images.
"""
import io
from typing import List
from fastapi import UploadFile
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image


async def extract_text_from_pdf(file: UploadFile) -> str:
    """Extract text content from a PDF file."""
    try:
        content = await file.read()
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_parts = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        await file.seek(0)  # Reset file pointer
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PDF: {str(e)}")


async def extract_text_from_docx(file: UploadFile) -> str:
    """Extract text content from a Word document."""
    try:
        content = await file.read()
        doc_file = io.BytesIO(content)
        doc = Document(doc_file)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        await file.seek(0)  # Reset file pointer
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from DOCX: {str(e)}")


async def extract_text_from_pptx(file: UploadFile) -> str:
    """Extract text content from a PowerPoint presentation."""
    try:
        content = await file.read()
        ppt_file = io.BytesIO(content)
        presentation = Presentation(ppt_file)
        
        text_parts = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
        
        await file.seek(0)  # Reset file pointer
        return "\n\n".join(text_parts)
    except Exception as e:
        raise ValueError(f"Failed to extract text from PPTX: {str(e)}")


async def extract_text_from_image(file: UploadFile) -> str:
    """
    Extract text from an image file.
    Note: This returns the image data for Gemini to process with OCR.
    Actual OCR is handled by Gemini API.
    """
    try:
        content = await file.read()
        # Verify it's a valid image
        image = Image.open(io.BytesIO(content))
        
        await file.seek(0)  # Reset file pointer
        return f"[Image file: {file.filename}, Size: {image.size}, Format: {image.format}]"
    except Exception as e:
        raise ValueError(f"Failed to process image: {str(e)}")


async def process_multiple_files(files: List[UploadFile]) -> str:
    """
    Process multiple files and combine their content.
    Automatically detects file type and uses appropriate extraction method.
    """
    combined_content = []
    
    for file in files:
        filename = file.filename.lower()
        
        try:
            if filename.endswith('.pdf'):
                text = await extract_text_from_pdf(file)
                combined_content.append(f"=== Content from {file.filename} ===\n{text}")
            
            elif filename.endswith(('.docx', '.doc')):
                text = await extract_text_from_docx(file)
                combined_content.append(f"=== Content from {file.filename} ===\n{text}")
            
            elif filename.endswith(('.pptx', '.ppt')):
                text = await extract_text_from_pptx(file)
                combined_content.append(f"=== Content from {file.filename} ===\n{text}")
            
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                text = await extract_text_from_image(file)
                combined_content.append(f"=== Image: {file.filename} ===\n{text}")
            
            else:
                # Try to read as plain text
                content = await file.read()
                text = content.decode('utf-8', errors='ignore')
                combined_content.append(f"=== Content from {file.filename} ===\n{text}")
                await file.seek(0)
        
        except Exception as e:
            combined_content.append(f"=== Error processing {file.filename} ===\n{str(e)}")
    
    return "\n\n".join(combined_content)


def get_supported_file_types() -> List[str]:
    """Return list of supported file extensions."""
    return [
        '.pdf',
        '.docx', '.doc',
        '.pptx', '.ppt',
        '.png', '.jpg', '.jpeg', '.gif', '.bmp',
        '.txt', '.md'
    ]

